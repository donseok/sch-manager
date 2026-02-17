"""
42병동 간호사 근무표 관리 시스템 - 퀵 매뉴얼 생성 스크립트
실행: py -3 docs/generate_quick_manual.py
출력: docs/quick-manual.pptx, docs/quick-manual.pdf
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu
from reportlab.lib.pagesizes import A4, landscape
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfgen import canvas
from reportlab.lib.colors import HexColor


OUTPUT_DIR = Path(__file__).parent
PPTX_PATH = OUTPUT_DIR / "quick-manual.pptx"
PDF_PATH = OUTPUT_DIR / "quick-manual.pdf"

FONT_NAME = "MalgunGothic"
FONT_BOLD = "MalgunGothicBold"
FONT_PATH = Path(r"C:\Windows\Fonts\malgun.ttf")
FONT_BOLD_PATH = Path(r"C:\Windows\Fonts\malgunbd.ttf")

# Design tokens
NAVY = RGBColor(0x1E, 0x3A, 0x5F)
EMERALD = RGBColor(0x10, 0xB9, 0x81)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1F, 0x2A, 0x40)
MID_TEXT = RGBColor(0x4B, 0x55, 0x63)
LIGHT_BG = RGBColor(0xF0, 0xF4, 0xF8)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
EMERALD_LIGHT = RGBColor(0xEC, 0xFD, 0xF5)
NAVY_LIGHT = RGBColor(0xEF, 0xF2, 0xF7)
BORDER_COLOR = RGBColor(0xE2, 0xE8, 0xF0)

# PDF equivalents
PDF_NAVY = HexColor("#1E3A5F")
PDF_EMERALD = HexColor("#10B981")
PDF_DARK = HexColor("#1F2A40")
PDF_MID = HexColor("#4B5563")
PDF_LIGHT_BG = HexColor("#F0F4F8")
PDF_CARD_BG = HexColor("#FFFFFF")
PDF_EMERALD_LIGHT = HexColor("#ECFDF5")
PDF_NAVY_LIGHT = HexColor("#EFF2F7")
PDF_BORDER = HexColor("#E2E8F0")


# ─── PPTX Helpers ───────────────────────────────────────────────────────────

def add_rounded_rect(slide, left, top, width, height, fill_rgb, border_rgb=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if border_rgb:
        shape.line.color.rgb = border_rgb
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    # Adjust corner radius
    shape.adjustments[0] = 0.04
    return shape


def add_text(slide, left, top, width, height, text, font_size=12,
             bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return box


def add_multiline_text(slide, left, top, width, height, lines, font_size=11,
                       color=DARK_TEXT, line_spacing=1.2):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(2)
        p.line_spacing = line_spacing
    return box


def add_header_bar(slide, text):
    """Top navy header bar with title"""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.85),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    bar.shadow.inherit = False

    # Emerald accent line
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(0.85), Inches(13.333), Inches(0.04),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = EMERALD
    accent.line.fill.background()
    accent.shadow.inherit = False

    add_text(slide, Inches(0.5), Inches(0.15), Inches(9), Inches(0.55),
             text, font_size=22, bold=True, color=WHITE)


def add_section_card(slide, left, top, width, height, number, title, items,
                     accent_color=EMERALD):
    """Card with number badge + title + bullet items"""
    card = add_rounded_rect(slide, left, top, width, height, CARD_BG, BORDER_COLOR)

    # Number badge
    badge_size = Inches(0.35)
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left + Inches(0.2), top + Inches(0.18), badge_size, badge_size,
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = accent_color
    badge.line.fill.background()
    badge.shadow.inherit = False
    badge.adjustments[0] = 0.15
    # Number text inside badge
    btf = badge.text_frame
    btf.clear()
    bp = btf.paragraphs[0]
    bp.text = str(number)
    bp.font.name = "Malgun Gothic"
    bp.font.size = Pt(14)
    bp.font.bold = True
    bp.font.color.rgb = WHITE
    bp.alignment = PP_ALIGN.CENTER
    btf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Title next to badge
    add_text(slide, left + Inches(0.65), top + Inches(0.15),
             width - Inches(0.85), Inches(0.35),
             title, font_size=14, bold=True, color=NAVY)

    # Bullet items
    bullet_lines = [f"  •  {item}" for item in items]
    add_multiline_text(
        slide, left + Inches(0.25), top + Inches(0.55),
        width - Inches(0.45), height - Inches(0.6),
        bullet_lines, font_size=10.5, color=MID_TEXT,
    )


def add_code_table(slide, left, top, width, height, codes):
    """Colored code reference strip"""
    card = add_rounded_rect(slide, left, top, width, height, NAVY_LIGHT, BORDER_COLOR)
    n = len(codes)
    cell_w = (width - Inches(0.3)) / n
    for i, (code, label, c_rgb) in enumerate(codes):
        x = left + Inches(0.15) + cell_w * i
        # Code letter
        add_text(slide, x, top + Inches(0.08), cell_w, Inches(0.3),
                 code, font_size=16, bold=True, color=c_rgb,
                 alignment=PP_ALIGN.CENTER)
        # Label
        add_text(slide, x, top + Inches(0.35), cell_w, Inches(0.22),
                 label, font_size=9, bold=False, color=MID_TEXT,
                 alignment=PP_ALIGN.CENTER)


# ─── Page 1: 시스템 개요 + 근무표 작성/편집 ─────────────────────────────────

def build_page1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = LIGHT_BG

    add_header_bar(slide,
                   "42병동 간호사 근무표 관리 시스템  |  퀵 매뉴얼 (1/2)")

    # Subtitle
    add_text(slide, Inches(0.5), Inches(1.05), Inches(12), Inches(0.35),
             "시스템 개요 및 근무표 작성 · 편집 워크플로우",
             font_size=13, bold=False, color=MID_TEXT)

    # ── Row 1: 3 cards ──
    card_w = Inches(4.0)
    card_h = Inches(2.2)
    gap = Inches(0.33)
    start_x = Inches(0.5)
    row1_y = Inches(1.5)

    add_section_card(slide, start_x, row1_y, card_w, card_h,
                     1, "시스템 접속 및 메뉴",
                     [
                         "브라우저에서 시스템 URL 접속",
                         "좌측 사이드바 메뉴 구성:",
                         "  📊 대시보드 — 금일 현황 확인",
                         "  📋 근무표 관리 — 생성/편집/확정",
                         "  👤 간호사 관리 — 등록/수정/삭제",
                     ], accent_color=NAVY)

    add_section_card(slide, start_x + card_w + gap, row1_y, card_w, card_h,
                     2, "근무표 생성 → 편집",
                     [
                         "근무표 관리 → [새 근무표] 클릭",
                         "연도/월 선택 후 생성",
                         "수간호사·책임간호사 자동 배정 (평일=D, 주말=O)",
                         "셀 클릭하여 근무코드 변경",
                     ], accent_color=EMERALD)

    add_section_card(slide, start_x + (card_w + gap) * 2, row1_y, card_w, card_h,
                     3, "저장 → 확정",
                     [
                         "[저장] 클릭 → 변경사항 DB 반영",
                         "모든 간호사의 모든 날짜 입력 완료 필요",
                         "[확정] 클릭 → 월별 1개만 확정 가능",
                         "확정 후 대시보드에 자동 반영",
                     ], accent_color=NAVY)

    # ── Row 2: Input methods card + Code reference ──
    row2_y = Inches(3.95)

    # Input methods card (wider)
    input_card_w = Inches(7.6)
    input_card_h = Inches(2.7)
    add_section_card(slide, start_x, row2_y, input_card_w, input_card_h,
                     4, "셀 입력 방법",
                     [
                         "클릭 입력: 셀 클릭 → 팝업에서 근무코드 선택",
                         "드래그 선택: 마우스 드래그로 여러 셀 범위 선택 후 일괄 입력",
                         "Shift + 클릭: 시작 셀 클릭 → Shift + 끝 셀 클릭으로 범위 확장",
                         "Ctrl+C / Ctrl+V: 선택 영역 복사 → 다른 위치에 붙여넣기",
                         "탭 구분자(TSV) 형식 지원 — 엑셀에서 복사한 데이터 붙여넣기 가능",
                     ], accent_color=EMERALD)

    # Code reference
    code_card_x = start_x + input_card_w + gap
    code_card_w = Inches(4.73)

    add_text(slide, code_card_x + Inches(0.15), row2_y + Inches(0.05),
             code_card_w, Inches(0.3),
             "근무코드 안내", font_size=13, bold=True, color=NAVY)

    code_top = row2_y + Inches(0.4)
    codes = [
        ("D", "주간", RGBColor(0x25, 0x63, 0xEB)),
        ("E", "저녁", RGBColor(0xD9, 0x77, 0x06)),
        ("N", "야간", RGBColor(0x7C, 0x3A, 0xED)),
        ("O", "공휴", RGBColor(0x10, 0xB9, 0x81)),
    ]
    add_code_table(slide, code_card_x, code_top, code_card_w, Inches(0.65), codes)

    codes2 = [
        ("X", "휴무", RGBColor(0x64, 0x74, 0x8B)),
        ("T", "교육", RGBColor(0xEC, 0x48, 0x99)),
        ("B", "기타", RGBColor(0x8B, 0x5C, 0xF6)),
    ]
    add_code_table(slide, code_card_x, code_top + Inches(0.78),
                   code_card_w, Inches(0.65), codes2)

    # Workflow summary strip at bottom
    strip_y = Inches(6.85)
    strip = add_rounded_rect(slide, start_x, strip_y,
                             Inches(12.33), Inches(0.45),
                             EMERALD_LIGHT, EMERALD)
    add_text(slide, start_x + Inches(0.3), strip_y + Inches(0.05),
             Inches(12), Inches(0.35),
             "워크플로우 요약:   간호사 등록  →  새 근무표 생성  →  셀 입력  →  저장  →  확정  →  대시보드 반영",
             font_size=11, bold=True, color=NAVY, alignment=PP_ALIGN.LEFT)


# ─── Page 2: 주요 기능 + 대시보드 활용 ───────────────────────────────────────

def build_page2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = LIGHT_BG

    add_header_bar(slide,
                   "42병동 간호사 근무표 관리 시스템  |  퀵 매뉴얼 (2/2)")

    add_text(slide, Inches(0.5), Inches(1.05), Inches(12), Inches(0.35),
             "주요 기능 및 대시보드 활용",
             font_size=13, bold=False, color=MID_TEXT)

    start_x = Inches(0.5)
    gap = Inches(0.33)

    # ── Row 1: 3 cards ──
    card_w = Inches(4.0)
    card_h = Inches(2.1)
    row1_y = Inches(1.5)

    add_section_card(slide, start_x, row1_y, card_w, card_h,
                     5, "엑셀 다운로드 / 인쇄",
                     [
                         "[엑셀 다운로드] → .xlsx 파일 저장",
                         "근무표 서식 그대로 다운로드",
                         "[인쇄] → 인쇄 전용 레이아웃으로 출력",
                         "브라우저 인쇄 기능 활용",
                     ], accent_color=NAVY)

    add_section_card(slide, start_x + card_w + gap, row1_y, card_w, card_h,
                     6, "이전 월 참조 / 사원 추가",
                     [
                         "[이전 월 불러오기] → 전월 확정 데이터 참조",
                         "[사원 추가] → 근무표에 간호사 추가",
                         "[사원 제거] → 근무표에서 간호사 제외",
                         "[초기화] → 전체 셀 초기화 (저장 전 확인!)",
                     ], accent_color=EMERALD)

    add_section_card(slide, start_x + (card_w + gap) * 2, row1_y, card_w, card_h,
                     7, "확정 / 확정 취소",
                     [
                         "확정 조건: 저장 완료 + 빈 셀 없음",
                         "월별 1개 근무표만 확정 가능",
                         "확정 후 편집/저장/초기화 잠김",
                         "[확정 취소] → 다시 수정 가능",
                     ], accent_color=NAVY)

    # ── Row 2: Dashboard card + Nurse management card ──
    row2_y = Inches(3.85)
    dash_w = Inches(7.6)
    dash_h = Inches(2.8)

    add_section_card(slide, start_x, row2_y, dash_w, dash_h,
                     8, "대시보드 활용 (확정된 근무표 기반)",
                     [
                         "금일 근무현황 — 오늘 D/E/N/O/X별 간호사 목록 확인",
                         "주간 미리보기 — 이번 주 근무 배치 한눈에 파악",
                         "공정성 지표 — 간호사별 D/E/N/O/X 횟수 분포 확인",
                         "알림 확인 — 연속 근무 5일 이상, N→D 연속 배치 등 경고",
                         "인력 현황 — 직위별 간호사 수, 근무 중/휴무 인원 통계",
                         "※ 대시보드는 확정된 근무표만 표시됩니다",
                     ], accent_color=EMERALD)

    # Nurse management card
    nurse_x = start_x + dash_w + gap
    nurse_w = Inches(4.73)
    nurse_h = Inches(2.8)

    add_section_card(slide, nurse_x, row2_y, nurse_w, nurse_h,
                     9, "간호사 관리",
                     [
                         "간호사 관리 메뉴에서 등록/수정/삭제",
                         "필수 정보: 사원번호, 사원명, 직위",
                         "직위: HN(수간호사), CN(책임간호사),",
                         "        AN(주임간호사), RN(일반간호사)",
                         "정렬순서로 근무표 표시 순서 결정",
                         "팀 배정 및 활성화 상태 관리 가능",
                     ], accent_color=NAVY)

    # Tips strip at bottom
    strip_y = Inches(6.85)
    strip = add_rounded_rect(slide, start_x, strip_y,
                             Inches(12.33), Inches(0.45),
                             EMERALD_LIGHT, EMERALD)
    add_text(slide, start_x + Inches(0.3), strip_y + Inches(0.05),
             Inches(12), Inches(0.35),
             "💡 팁:  확정 전에는 자유롭게 수정 가능  |  변경 이력에서 수정 내역 추적 가능  |  문의: 시스템 관리자에게 연락",
             font_size=11, bold=True, color=NAVY, alignment=PP_ALIGN.LEFT)


# ─── PPTX Build ──────────────────────────────────────────────────────────────

def make_pptx() -> None:
    prs = Presentation()
    # Set widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_page1(prs)
    build_page2(prs)
    prs.save(PPTX_PATH)
    print(f"[OK] PPTX created: {PPTX_PATH}")


# ─── PDF Build (ReportLab) ──────────────────────────────────────────────────

def wrap_lines(text: str, font: str, size: int, max_width: float) -> list[str]:
    words = list(text)  # character-level for Korean
    lines: list[str] = []
    current = ""
    for ch in words:
        candidate = current + ch
        if pdfmetrics.stringWidth(candidate, font, size) <= max_width:
            current = candidate
        else:
            if current:
                lines.append(current)
            current = ch
    if current:
        lines.append(current)
    return lines


def draw_rounded_rect(c, x, y, w, h, r, fill=None, stroke=None, stroke_w=0.5):
    """Draw rounded rectangle. y is bottom-left."""
    c.saveState()
    if fill:
        c.setFillColor(fill)
    if stroke:
        c.setStrokeColor(stroke)
        c.setLineWidth(stroke_w)
    p = c.beginPath()
    p.roundRect(x, y, w, h, r)
    if fill and stroke:
        c.drawPath(p, fill=1, stroke=1)
    elif fill:
        c.drawPath(p, fill=1, stroke=0)
    else:
        c.drawPath(p, fill=0, stroke=1)
    c.restoreState()


def draw_header_bar(c, pw, ph, title):
    # Navy bar
    c.setFillColor(PDF_NAVY)
    c.rect(0, ph - 52, pw, 52, fill=1, stroke=0)
    # Emerald accent line
    c.setFillColor(PDF_EMERALD)
    c.rect(0, ph - 55, pw, 3, fill=1, stroke=0)
    # Title
    c.setFillColor(HexColor("#FFFFFF"))
    c.setFont(FONT_BOLD, 16)
    c.drawString(28, ph - 38, title)


def draw_section_card(c, x, y, w, h, number, title, items, accent_hex="#10B981"):
    """y = top of card (from page top), converts internally"""
    accent = HexColor(accent_hex)

    # Card background
    draw_rounded_rect(c, x, y - h, w, h, 6, fill=PDF_CARD_BG, stroke=PDF_BORDER)

    # Number badge
    badge_size = 20
    bx = x + 10
    by = y - 24
    draw_rounded_rect(c, bx, by, badge_size, badge_size, 4, fill=accent)
    c.setFillColor(HexColor("#FFFFFF"))
    c.setFont(FONT_BOLD, 11)
    nw = pdfmetrics.stringWidth(str(number), FONT_BOLD, 11)
    c.drawString(bx + (badge_size - nw) / 2, by + 5, str(number))

    # Title
    c.setFillColor(PDF_NAVY)
    c.setFont(FONT_BOLD, 11)
    c.drawString(x + 36, by + 4, title)

    # Items
    c.setFont(FONT_NAME, 8)
    c.setFillColor(PDF_MID)
    iy = y - 42
    for item in items:
        text = f"•  {item}"
        lines = wrap_lines(text, FONT_NAME, 8, w - 24)
        for line in lines:
            c.drawString(x + 14, iy, line)
            iy -= 13
        iy -= 1


def draw_code_strip(c, x, y, w, codes):
    """Draw horizontal code reference strip"""
    strip_h = 36
    draw_rounded_rect(c, x, y - strip_h, w, strip_h, 4, fill=PDF_NAVY_LIGHT, stroke=PDF_BORDER)
    n = len(codes)
    cell_w = (w - 10) / n
    for i, (code, label, hex_color) in enumerate(codes):
        cx = x + 5 + cell_w * i + cell_w / 2
        c.setFillColor(HexColor(hex_color))
        c.setFont(FONT_BOLD, 12)
        cw = pdfmetrics.stringWidth(code, FONT_BOLD, 12)
        c.drawString(cx - cw / 2, y - 15, code)
        c.setFillColor(PDF_MID)
        c.setFont(FONT_NAME, 7)
        lw = pdfmetrics.stringWidth(label, FONT_NAME, 7)
        c.drawString(cx - lw / 2, y - 28, label)


def draw_bottom_strip(c, x, y, w, text):
    strip_h = 24
    draw_rounded_rect(c, x, y, w, strip_h, 4,
                      fill=PDF_EMERALD_LIGHT, stroke=PDF_EMERALD, stroke_w=0.75)
    c.setFillColor(PDF_NAVY)
    c.setFont(FONT_BOLD, 8)
    c.drawString(x + 12, y + 7, text)


def build_pdf_page1(c, pw, ph):
    # Background
    c.setFillColor(PDF_LIGHT_BG)
    c.rect(0, 0, pw, ph, fill=1, stroke=0)

    draw_header_bar(c, pw, ph,
                    "42병동 간호사 근무표 관리 시스템  |  퀵 매뉴얼 (1/2)")

    # Subtitle
    c.setFillColor(PDF_MID)
    c.setFont(FONT_NAME, 10)
    c.drawString(28, ph - 74, "시스템 개요 및 근무표 작성 · 편집 워크플로우")

    margin = 28
    gap = 14
    card_w = (pw - 2 * margin - 2 * gap) / 3
    row1_top = ph - 90

    card_h = 130
    draw_section_card(c, margin, row1_top, card_w, card_h,
                      1, "시스템 접속 및 메뉴",
                      [
                          "브라우저에서 시스템 URL 접속",
                          "좌측 사이드바 메뉴 구성:",
                          "  대시보드 — 금일 현황 확인",
                          "  근무표 관리 — 생성/편집/확정",
                          "  간호사 관리 — 등록/수정/삭제",
                      ], accent_hex="#1E3A5F")

    draw_section_card(c, margin + card_w + gap, row1_top, card_w, card_h,
                      2, "근무표 생성 → 편집",
                      [
                          "근무표 관리 → [새 근무표] 클릭",
                          "연도/월 선택 후 생성",
                          "수간호사·책임간호사 자동 배정",
                          "  (평일=D, 주말=O)",
                          "셀 클릭하여 근무코드 변경",
                      ], accent_hex="#10B981")

    draw_section_card(c, margin + 2 * (card_w + gap), row1_top, card_w, card_h,
                      3, "저장 → 확정",
                      [
                          "[저장] 클릭 → 변경사항 DB 반영",
                          "모든 간호사의 날짜 입력 완료 필요",
                          "[확정] 클릭 → 월별 1개만 확정 가능",
                          "확정 후 대시보드에 자동 반영",
                      ], accent_hex="#1E3A5F")

    # Row 2
    row2_top = row1_top - card_h - gap
    input_w = pw * 0.6 - margin
    input_h = 155

    draw_section_card(c, margin, row2_top, input_w, input_h,
                      4, "셀 입력 방법",
                      [
                          "클릭 입력: 셀 클릭 → 팝업에서 근무코드 선택",
                          "드래그 선택: 마우스 드래그로 여러 셀 범위 선택",
                          "Shift + 클릭: 시작셀 → Shift + 끝셀로 범위 확장",
                          "Ctrl+C / Ctrl+V: 선택 영역 복사/붙여넣기",
                          "탭 구분자(TSV) — 엑셀 데이터 붙여넣기 가능",
                      ], accent_hex="#10B981")

    # Code reference area
    code_x = margin + input_w + gap
    code_w = pw - code_x - margin

    c.setFillColor(PDF_NAVY)
    c.setFont(FONT_BOLD, 10)
    c.drawString(code_x + 6, row2_top - 14, "근무코드 안내")

    codes1 = [
        ("D", "주간", "#2563EB"), ("E", "저녁", "#D97706"),
        ("N", "야간", "#7C3AED"), ("O", "공휴", "#10B981"),
    ]
    draw_code_strip(c, code_x, row2_top - 24, code_w, codes1)

    codes2 = [
        ("X", "휴무", "#64748B"), ("T", "교육", "#EC4899"),
        ("B", "기타", "#8B5CF6"),
    ]
    draw_code_strip(c, code_x, row2_top - 68, code_w, codes2)

    # Bottom strip
    draw_bottom_strip(c, margin, 18, pw - 2 * margin,
                      "워크플로우 요약:   간호사 등록  →  새 근무표 생성  →  셀 입력  →  저장  →  확정  →  대시보드 반영")


def build_pdf_page2(c, pw, ph):
    c.setFillColor(PDF_LIGHT_BG)
    c.rect(0, 0, pw, ph, fill=1, stroke=0)

    draw_header_bar(c, pw, ph,
                    "42병동 간호사 근무표 관리 시스템  |  퀵 매뉴얼 (2/2)")

    c.setFillColor(PDF_MID)
    c.setFont(FONT_NAME, 10)
    c.drawString(28, ph - 74, "주요 기능 및 대시보드 활용")

    margin = 28
    gap = 14
    card_w = (pw - 2 * margin - 2 * gap) / 3
    row1_top = ph - 90
    card_h = 120

    draw_section_card(c, margin, row1_top, card_w, card_h,
                      5, "엑셀 다운로드 / 인쇄",
                      [
                          "[엑셀 다운로드] → .xlsx 파일 저장",
                          "근무표 서식 그대로 다운로드",
                          "[인쇄] → 인쇄 전용 레이아웃 출력",
                          "브라우저 인쇄 기능 활용",
                      ], accent_hex="#1E3A5F")

    draw_section_card(c, margin + card_w + gap, row1_top, card_w, card_h,
                      6, "이전 월 참조 / 사원 추가",
                      [
                          "[이전 월 불러오기] → 전월 데이터 참조",
                          "[사원 추가] → 근무표에 간호사 추가",
                          "[사원 제거] → 근무표에서 간호사 제외",
                          "[초기화] → 전체 셀 초기화",
                      ], accent_hex="#10B981")

    draw_section_card(c, margin + 2 * (card_w + gap), row1_top, card_w, card_h,
                      7, "확정 / 확정 취소",
                      [
                          "확정 조건: 저장 완료 + 빈 셀 없음",
                          "월별 1개 근무표만 확정 가능",
                          "확정 후 편집/저장/초기화 잠김",
                          "[확정 취소] → 다시 수정 가능",
                      ], accent_hex="#1E3A5F")

    # Row 2
    row2_top = row1_top - card_h - gap
    dash_w = pw * 0.6 - margin
    dash_h = 165

    draw_section_card(c, margin, row2_top, dash_w, dash_h,
                      8, "대시보드 활용 (확정된 근무표 기반)",
                      [
                          "금일 근무현황 — D/E/N/O/X별 간호사 목록",
                          "주간 미리보기 — 이번 주 근무 배치 파악",
                          "공정성 지표 — 간호사별 근무 횟수 분포",
                          "알림 — 연속근무 5일↑, N→D 경고",
                          "인력 현황 — 직위별/근무상태별 통계",
                          "※ 확정된 근무표만 표시됩니다",
                      ], accent_hex="#10B981")

    nurse_x = margin + dash_w + gap
    nurse_w = pw - nurse_x - margin
    nurse_h = 165

    draw_section_card(c, nurse_x, row2_top, nurse_w, nurse_h,
                      9, "간호사 관리",
                      [
                          "간호사 관리 메뉴에서 등록/수정/삭제",
                          "필수 정보: 사원번호, 사원명, 직위",
                          "직위: HN(수간호사), CN(책임간호사),",
                          "       AN(주임간호사), RN(일반간호사)",
                          "정렬순서로 근무표 표시 순서 결정",
                          "팀 배정 및 활성화 상태 관리",
                      ], accent_hex="#1E3A5F")

    draw_bottom_strip(c, margin, 18, pw - 2 * margin,
                      "팁:  확정 전에는 자유롭게 수정 가능  |  변경 이력에서 수정 내역 추적 가능  |  문의: 시스템 관리자에게 연락")


def make_pdf() -> None:
    if not FONT_PATH.exists():
        raise FileNotFoundError(f"Font not found: {FONT_PATH}")

    pdfmetrics.registerFont(TTFont(FONT_NAME, str(FONT_PATH)))
    if FONT_BOLD_PATH.exists():
        pdfmetrics.registerFont(TTFont(FONT_BOLD, str(FONT_BOLD_PATH)))
    else:
        pdfmetrics.registerFont(TTFont(FONT_BOLD, str(FONT_PATH)))

    page_w, page_h = landscape(A4)
    c_obj = canvas.Canvas(str(PDF_PATH), pagesize=landscape(A4))

    build_pdf_page1(c_obj, page_w, page_h)
    c_obj.showPage()
    build_pdf_page2(c_obj, page_w, page_h)
    c_obj.showPage()

    c_obj.save()
    print(f"[OK] PDF created:  {PDF_PATH}")


# ─── Main ────────────────────────────────────────────────────────────────────

def main() -> None:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    make_pptx()
    make_pdf()
    print("Done!")


if __name__ == "__main__":
    main()
